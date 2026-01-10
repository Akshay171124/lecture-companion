from __future__ import annotations

from io import BytesIO
from typing import Tuple
from datetime import datetime, timezone

from pypdf import PdfReader
from pptx import Presentation


def extract_pdf(data: bytes) -> str:
    reader = PdfReader(BytesIO(data))
    parts: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"\n--- page {i} ---\n{text}")
    return "\n".join(parts).strip()


def extract_pptx(data: bytes) -> str:
    prs = Presentation(BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    slide_text.append(t)
        if slide_text:
            parts.append(f"\n--- slide {i} ---\n" + "\n".join(slide_text))
    return "\n".join(parts).strip()


def extract_text(filename: str, mime_type: str | None, data: bytes) -> Tuple[str, str]:
    """
    Returns: (status, extracted_text_or_error)
    """
    name = filename.lower()

    try:
        if name.endswith(".pdf") or (mime_type == "application/pdf"):
            text = extract_pdf(data)
            return ("EXTRACTED", text)

        if name.endswith(".pptx") or (mime_type in {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}):
            text = extract_pptx(data)
            return ("EXTRACTED", text)

        return ("FAILED", "Unsupported file type. Only PDF and PPTX are supported in Phase 1.1.")
    except Exception as e:
        return ("FAILED", f"Extraction error: {e}")
